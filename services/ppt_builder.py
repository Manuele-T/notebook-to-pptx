import logging
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from io import BytesIO

logger = logging.getLogger(__name__)


def generate_pptx(slides_data: list[dict]) -> BytesIO:
    """
    Rebuilds each slide in PPTX format with proper text fitting and image placement.
    """
    logger.info(f"Starting PPTX generation with {len(slides_data)} slides")
    prs = Presentation()

    for idx, slide_data in enumerate(slides_data):
        l_type = slide_data.get("layout_type", "title_and_content")
        title_text = slide_data.get("title", "")
        body_text_list = slide_data.get("body_text", [])
        figures = slide_data.get("figures", [])
        speaker_notes = slide_data.get("speaker_notes", "")
        
        logger.info(f"Building slide {idx+1}: {l_type}")
        
        # Always use blank layout for maximum control
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        has_figures = len(figures) > 0
        has_body = len(body_text_list) > 0
        
        # --- TITLE (always at top) ---
        if title_text:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), 
                Inches(9), Inches(0.8)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(28)
            p.font.bold = True
            p.alignment = PP_ALIGN.LEFT
        
        # --- CONTENT LAYOUT LOGIC ---
        if l_type == "title_only":
            if has_figures:
                _add_centered_image(slide, figures[0], top=1.5, max_width=8, max_height=5.5)
        
        elif l_type == "two_column":
            if has_figures and len(figures) >= 2:
                # Two figures side by side
                _add_image_at(slide, figures[0], left=0.5, top=1.5, max_width=4.2, max_height=4)
                _add_image_at(slide, figures[1], left=5.3, top=1.5, max_width=4.2, max_height=4)
                if has_body:
                    _add_body_text(slide, body_text_list, left=0.5, top=5.7, width=9, height=1.5, font_size=11)
            elif has_figures:
                _add_body_text(slide, body_text_list, left=0.5, top=1.3, width=4.5, height=5.5, font_size=12)
                _add_image_at(slide, figures[0], left=5.2, top=1.5, max_width=4.3, max_height=5)
            else:
                _add_body_text(slide, body_text_list, left=0.5, top=1.3, width=9, height=5.8, font_size=13)
        
        elif l_type == "diagram_heavy":
            if has_figures:
                _add_centered_image(slide, figures[0], top=1.3, max_width=9, max_height=5.2)
            if has_body:
                combined = " | ".join(body_text_list)
                _add_body_text(slide, [combined], left=0.5, top=6.5, width=9, height=0.8, font_size=10)
        
        else:  # title_and_content, mixed_freeform, or default
            if has_figures and has_body:
                _add_body_text(slide, body_text_list, left=0.5, top=1.3, width=5.5, height=5.8, font_size=13)
                _add_image_at(slide, figures[0], left=6.2, top=1.5, max_width=3.5, max_height=5)
            elif has_figures:
                _add_centered_image(slide, figures[0], top=1.5, max_width=8, max_height=5.5)
            elif has_body:
                _add_body_text(slide, body_text_list, left=0.5, top=1.3, width=9, height=5.8, font_size=14)
        
        # --- SPEAKER NOTES ---
        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    logger.info("PPTX generation complete")
    return pptx_buffer


def _add_body_text(slide, items: list[str], left: float, top: float, 
                   width: float, height: float, font_size: int = 14):
    """Add a text box with content that fits within bounds."""
    if not items:
        return
    
    text_box = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    tf = text_box.text_frame
    tf.word_wrap = True
    
    # Adjust font size based on content length
    total_chars = sum(len(item) for item in items)
    if total_chars > 1000:
        font_size = min(font_size, 10)
    elif total_chars > 600:
        font_size = min(font_size, 11)
    elif total_chars > 400:
        font_size = min(font_size, 12)
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = item.strip()
        p.font.size = Pt(font_size)
        p.space_after = Pt(4)
        p.space_before = Pt(2)


def _add_centered_image(slide, figure: dict, top: float, max_width: float, max_height: float):
    """Add a centered image that fits within bounds."""
    if "image_bytes" not in figure:
        return
    
    image_stream = BytesIO(figure["image_bytes"])
    left = (10 - max_width) / 2
    
    try:
        pic = slide.shapes.add_picture(image_stream, Inches(left), Inches(top), width=Inches(max_width))
        
        if pic.height.inches > max_height:
            scale = max_height / pic.height.inches
            pic.width = Emu(int(pic.width * scale))
            pic.height = Emu(int(pic.height * scale))
            pic.left = Inches((10 - pic.width.inches) / 2)
    except Exception as e:
        logger.error(f"Error adding image: {e}")


def _add_image_at(slide, figure: dict, left: float, top: float, max_width: float, max_height: float):
    """Add an image at a specific position."""
    if "image_bytes" not in figure:
        return
    
    image_stream = BytesIO(figure["image_bytes"])
    
    try:
        pic = slide.shapes.add_picture(image_stream, Inches(left), Inches(top), width=Inches(max_width))
        
        if pic.height.inches > max_height:
            scale = max_height / pic.height.inches
            pic.width = Emu(int(pic.width * scale))
            pic.height = Emu(int(pic.height * scale))
    except Exception as e:
        logger.error(f"Error adding image: {e}")
